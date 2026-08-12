"""Fixed, network-free MCP adapters for catalog wave 3.

The adapters expose opaque file identifiers instead of paths.  All input,
output and persistent-memory roots are selected by the trusted sidecar.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import hmac
import io
import json
import math
import os
import re
import stat
import subprocess
import time
import zipfile
from pathlib import Path
from pathlib import PurePosixPath
from typing import Annotated, Any, Literal
from urllib.parse import unquote, urlsplit
from xml.etree import ElementTree

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pydantic import Field

from .file_artifacts import WAVE18A_BUILDERS, WAVE18A_TOOL_NAMES
from .file_analysis import WAVE18B_BUILDERS, WAVE18B_TOOL_NAMES
from .file_code_index import WAVE20_BUILDERS, WAVE20_TOOL_NAMES
from .file_wave26 import WAVE26_BUILDERS, WAVE26_TOOL_NAMES


WORKSPACE_PATTERN = re.compile(r"mcpws_[0-9a-f]{32}")
FILE_ID_PATTERN = re.compile(r"mcpf_[0-9a-f]{24}")
SAFE_NAME = re.compile(r"[A-Za-z0-9\u4e00-\u9fff][A-Za-z0-9\u4e00-\u9fff._ -]{0,119}")
MAX_INLINE_CHARS = 240_000
MAX_EXCEL_ROWS = 10_000
MAX_EXCEL_COLUMNS = 200
MAX_OFFICE_INPUT_BYTES = 10 * 1024 * 1024
MAX_OFFICE_ENTRIES = 10_000
MAX_OFFICE_MEMBER_BYTES = 100 * 1024 * 1024
MAX_OFFICE_UNCOMPRESSED_BYTES = 250 * 1024 * 1024
MAX_OFFICE_COMPRESSION_RATIO = 100
MAX_OFFICE_XML_BYTES = 16 * 1024 * 1024
MAX_OFFICE_XML_NODES = 500_000
MAX_OFFICE_XML_DEPTH = 64
MAX_OFFICE_XML_ATTRIBUTES = 500_000
MAX_OFFICE_XML_TEXT_CHARS = 8_000_000
MAX_OFFICE_EXTRACTED_CHARS = 500_000
MAX_OFFICE_SECTION_CHARS = 20_000
MAX_OFFICE_WIRE_BYTES = 2 * 1024 * 1024
MAX_OFFICE_SECTIONS = 10_000
MAX_OUTPUT_RENDER_SPEC_BYTES = 2 * 1024 * 1024
MAX_OUTPUT_RENDERED_BYTES = 50 * 1024 * 1024
MAX_OUTPUT_RENDER_CHARS = 500_000
MAX_OUTPUT_RENDER_COLUMNS = 200
MAX_OUTPUT_RENDER_CELLS = 100_000
MAX_OUTPUT_RENDER_SHEETS = 20
MAX_OUTPUT_RENDER_SLIDES = 100
OFFICE_HANDOFF_MARKER_NAME = ".modelmirror-office-parser.json"
OFFICE_HANDOFF_MARKER_OWNER = "modelmirror.file_assets.office_sidecar.v1"
OFFICE_HANDOFF_MARKER_MAX_BYTES = 4 * 1024
OUTPUT_RENDER_MARKER_NAME = ".modelmirror-output-renderer.json"
OUTPUT_RENDER_MARKER_OWNER = "modelmirror.file_assets.output_renderer.v1"
OUTPUT_RENDER_FORMATS = {
    "pdf": (".pdf", "application/pdf"),
    "docx": (
        ".docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ),
    "xlsx": (
        ".xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ),
    "pptx": (
        ".pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ),
}
_OFFICE_FIXED_SOURCE_NAMES = {
    "docx": "source.docx",
    "pptx": "source.pptx",
}
_OFFICE_SHA256_PATTERN = re.compile(r"[0-9a-f]{64}")

_OFFICE_MAIN_PARTS = {
    "docx": (
        "word/document.xml",
        "document",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
    ),
    "pptx": (
        "ppt/presentation.xml",
        "presentation",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    ),
}
_OFFICE_ALLOWED_COMPRESSION = {zipfile.ZIP_STORED, zipfile.ZIP_DEFLATED}
_OFFICE_ACTIVE_RELATIONSHIP_SUFFIXES = (
    "/oleobject",
    "/package",
    "/control",
    "/attachedtemplate",
    "/afchunk",
    "/externallink",
    "/vbaproject",
)

FileId = Annotated[
    str,
    Field(
        description="从当前受控工作区选择文件，不接受本机路径或 URI。",
        json_schema_extra={"x-modelmirror-input": "workspace-file"},
    ),
]
ArtifactName = Annotated[
    str,
    Field(
        description="产物文件名；产物只会写入可清理目录。",
        json_schema_extra={"x-modelmirror-input": "artifact-name"},
    ),
]

READ_ONLY = ToolAnnotations(
    readOnlyHint=True,
    destructiveHint=False,
    idempotentHint=True,
    openWorldHint=False,
)
ARTIFACT_CREATE = ToolAnnotations(
    readOnlyHint=False,
    destructiveHint=False,
    idempotentHint=False,
    openWorldHint=False,
)
STATE_WRITE = ToolAnnotations(
    readOnlyHint=False,
    destructiveHint=True,
    idempotentHint=False,
    openWorldHint=False,
)


def opaque_file_id(workspace_id: str, relative_path: str) -> str:
    digest = hashlib.sha256(
        f"{workspace_id}:{relative_path}".encode("utf-8")
    ).hexdigest()[:24]
    return f"mcpf_{digest}"


class WorkspaceContext:
    def __init__(self) -> None:
        workspace_id = os.getenv("MCP_FILE_WORKSPACE_ID", "").strip()
        if not WORKSPACE_PATTERN.fullmatch(workspace_id):
            raise RuntimeError("受控工作区标识无效。")
        self.workspace_id = workspace_id
        self.input_root = self._workspace_root("MCP_FILE_INPUT_ROOT", "/inputs")
        self.output_root = self._workspace_root("MCP_FILE_OUTPUT_ROOT", "/outputs")
        self.memory_root = self._workspace_root("MCP_FILE_MEMORY_ROOT", "/memory")
        self.output_root.mkdir(parents=True, exist_ok=True)
        self.memory_root.mkdir(parents=True, exist_ok=True)
        self._files: dict[str, Path] = {}
        if self.input_root.exists():
            for path in sorted(self.input_root.rglob("*")):
                if path.is_symlink() or not path.is_file():
                    continue
                relative = path.relative_to(self.input_root).as_posix()
                self._files[opaque_file_id(workspace_id, relative)] = path

    def _workspace_root(self, env_name: str, default: str) -> Path:
        base = Path(os.getenv(env_name, default)).resolve()
        root = (base / self.workspace_id).resolve()
        if root.parent != base:
            raise RuntimeError("受控工作区路径越界。")
        return root

    def resolve_file(self, file_id: str) -> Path:
        if not FILE_ID_PATTERN.fullmatch(str(file_id or "")):
            raise ValueError("必须选择当前工作区中的文件。")
        path = self._files.get(file_id)
        if path is None or not path.is_file() or path.is_symlink():
            raise ValueError("所选工作区文件不存在。")
        resolved = path.resolve()
        if self.input_root not in resolved.parents:
            raise ValueError("所选文件越过受控工作区。")
        return resolved

    def workspace_files(self) -> tuple[tuple[str, str, Path], ...]:
        """Return contained sealed-workspace files without exposing host paths."""

        values: list[tuple[str, str, Path]] = []
        for file_id, path in self._files.items():
            if path.is_symlink() or not path.is_file():
                raise ValueError("受控工作区文件集合已经变化。")
            resolved = path.resolve()
            if self.input_root not in resolved.parents:
                raise ValueError("受控工作区文件越界。")
            relative = resolved.relative_to(self.input_root).as_posix()
            values.append((file_id, relative, resolved))
        return tuple(sorted(values, key=lambda item: item[1]))

    def artifact_path(self, name: str, suffix: str) -> Path:
        clean = Path(str(name or "")).name
        if not SAFE_NAME.fullmatch(clean):
            raise ValueError("产物名称只能包含常用文字、数字、空格、点、下划线和连字符。")
        if not clean.lower().endswith(suffix.lower()):
            clean += suffix
        target = (self.output_root / clean).resolve()
        if target.parent != self.output_root or target.is_symlink():
            raise ValueError("产物路径无效。")
        return target

    @staticmethod
    def artifact_payload(path: Path) -> dict[str, Any]:
        data = path.read_bytes()
        return {
            "artifact_name": path.name,
            "relative_path": path.name,
            "size_bytes": len(data),
            "sha256": hashlib.sha256(data).hexdigest(),
        }


def _safe_note_path(root: Path, value: str, *, allow_missing: bool = False) -> Path:
    clean = str(value or "").strip().replace("\\", "/").strip("/")
    if clean.endswith(".md"):
        clean = clean[:-3]
    if not clean or len(clean) > 240:
        raise ValueError("笔记标识不能为空且不能超过 240 个字符。")
    parts = clean.split("/")
    if any(part in {"", ".", ".."} or not SAFE_NAME.fullmatch(part) for part in parts):
        raise ValueError("笔记路径包含不安全字符。")
    target = (root / ("/".join(parts) + ".md")).resolve()
    if root not in target.parents or target.is_symlink():
        raise ValueError("笔记路径越过持久工作区。")
    if not allow_missing and not target.is_file():
        raise ValueError("笔记不存在。")
    return target


def _note_payload(path: Path, root: Path, content: str) -> dict[str, Any]:
    return {
        "note": path.relative_to(root).as_posix(),
        "content": content[:MAX_INLINE_CHARS],
        "truncated": len(content) > MAX_INLINE_CHARS,
        "updated_at": path.stat().st_mtime,
    }


def build_basic_memory(context: WorkspaceContext) -> FastMCP:
    mcp = FastMCP("ModelMirror Basic Memory")
    notes = context.memory_root / "notes"
    notes.mkdir(parents=True, exist_ok=True)

    def read(value: str) -> dict[str, Any]:
        path = _safe_note_path(notes, value)
        return _note_payload(path, notes, path.read_text(encoding="utf-8"))

    @mcp.tool(annotations=READ_ONLY)
    def read_note(note: str) -> dict[str, Any]:
        """读取当前持久记忆库中的 Markdown 笔记。"""
        return read(note)

    @mcp.tool(annotations=READ_ONLY)
    def read_content(note: str) -> dict[str, Any]:
        """读取当前持久记忆库中的 Markdown 内容。"""
        return read(note)

    @mcp.tool(annotations=READ_ONLY)
    def view_note(note: str) -> dict[str, Any]:
        """查看当前持久记忆库中的笔记。"""
        return read(note)

    @mcp.tool(annotations=READ_ONLY)
    def search_notes(query: str, limit: int = 20) -> dict[str, Any]:
        """按标题和正文关键词搜索本地笔记。"""
        clean = str(query or "").strip().casefold()
        if not clean or len(clean) > 500:
            raise ValueError("搜索词必须包含 1 到 500 个字符。")
        matches: list[dict[str, Any]] = []
        for path in sorted(notes.rglob("*.md")):
            if path.is_symlink():
                continue
            text = path.read_text(encoding="utf-8", errors="replace")
            if clean in path.stem.casefold() or clean in text.casefold():
                position = text.casefold().find(clean)
                start = max(0, position - 120)
                matches.append({
                    "note": path.relative_to(notes).as_posix(),
                    "preview": text[start:start + 360],
                    "updated_at": path.stat().st_mtime,
                })
            if len(matches) >= max(1, min(int(limit), 100)):
                break
        return {"query": query, "matches": matches, "count": len(matches)}

    @mcp.tool(annotations=READ_ONLY)
    def search(query: str, limit: int = 20) -> dict[str, Any]:
        """兼容搜索入口，仅检索当前本地记忆库。"""
        return search_notes(query, limit)

    @mcp.tool(annotations=READ_ONLY)
    def fetch(note: str) -> dict[str, Any]:
        """兼容读取入口，仅访问当前本地记忆库。"""
        return read(note)

    @mcp.tool(annotations=READ_ONLY)
    def recent_activity(limit: int = 20) -> dict[str, Any]:
        """列出最近更新的本地笔记。"""
        paths = sorted(
            (path for path in notes.rglob("*.md") if not path.is_symlink()),
            key=lambda path: path.stat().st_mtime,
            reverse=True,
        )[:max(1, min(int(limit), 100))]
        return {"items": [
            {"note": path.relative_to(notes).as_posix(), "updated_at": path.stat().st_mtime}
            for path in paths
        ]}

    @mcp.tool(annotations=READ_ONLY)
    def list_directory(folder: str = "") -> dict[str, Any]:
        """列出记忆库中的 Markdown 笔记，不接受宿主目录。"""
        clean = str(folder or "").strip().replace("\\", "/").strip("/")
        target = notes if not clean else (notes / clean).resolve()
        if target != notes and notes not in target.parents:
            raise ValueError("目录越过记忆库。")
        if not target.is_dir() or target.is_symlink():
            raise ValueError("记忆目录不存在。")
        return {"folder": clean, "items": [
            {
                "name": path.name,
                "kind": "directory" if path.is_dir() else "note",
            }
            for path in sorted(target.iterdir())
            if not path.is_symlink() and (path.is_dir() or path.suffix == ".md")
        ][:500]}

    @mcp.tool(annotations=READ_ONLY)
    def build_context(note: str, depth: int = 1) -> dict[str, Any]:
        """读取笔记及最多两层 Wiki 链接上下文。"""
        depth = max(0, min(int(depth), 2))
        first = read(note)
        related: list[dict[str, Any]] = []
        frontier = re.findall(r"\[\[([^\]]+)\]\]", str(first["content"]))[:30]
        seen = {str(first["note"]).casefold()}
        for _ in range(depth):
            next_frontier: list[str] = []
            for name in frontier:
                try:
                    payload = read(name)
                except ValueError:
                    continue
                key = str(payload["note"]).casefold()
                if key in seen:
                    continue
                seen.add(key)
                related.append(payload)
                next_frontier.extend(re.findall(r"\[\[([^\]]+)\]\]", str(payload["content"]))[:10])
            frontier = next_frontier[:30]
        return {"root": first, "related": related[:50]}

    @mcp.tool(annotations=READ_ONLY)
    def basic_memory_diagnostics() -> dict[str, Any]:
        """返回本地、断网记忆库的非敏感状态。"""
        paths = [path for path in notes.rglob("*.md") if not path.is_symlink()]
        return {
            "mode": "local-only",
            "network": "disabled",
            "semantic_downloads": "disabled",
            "notes": len(paths),
        }

    @mcp.tool(annotations=STATE_WRITE)
    def write_note(title: str, content: str, folder: str = "") -> dict[str, Any]:
        """经目录审批后创建新的持久 Markdown 笔记。"""
        if not str(content).strip() or len(str(content).encode("utf-8")) > 1024 * 1024:
            raise ValueError("笔记正文不能为空且不能超过 1 MiB。")
        relative = "/".join(value for value in (str(folder).strip("/"), str(title).strip()) if value)
        path = _safe_note_path(notes, relative, allow_missing=True)
        if path.exists():
            raise ValueError("目标笔记已经存在，拒绝覆盖。")
        path.parent.mkdir(parents=True, exist_ok=True)
        body = f"---\ntitle: {title}\n---\n\n# {title}\n\n{content.strip()}\n"
        path.write_text(body, encoding="utf-8")
        return _note_payload(path, notes, body)

    @mcp.tool(annotations=STATE_WRITE)
    def edit_note(
        note: str,
        content: str,
        operation: Literal["append", "prepend", "replace"] = "append",
    ) -> dict[str, Any]:
        """经目录审批后修改已有持久笔记。"""
        path = _safe_note_path(notes, note)
        old = path.read_text(encoding="utf-8")
        if len(str(content).encode("utf-8")) > 1024 * 1024:
            raise ValueError("单次编辑内容不能超过 1 MiB。")
        if operation == "append":
            value = old.rstrip() + "\n\n" + str(content).strip() + "\n"
        elif operation == "prepend":
            value = str(content).strip() + "\n\n" + old
        else:
            value = str(content)
        path.write_text(value, encoding="utf-8")
        return _note_payload(path, notes, value)

    @mcp.tool(annotations=STATE_WRITE)
    def move_note(note: str, destination_folder: str) -> dict[str, Any]:
        """经目录审批后在当前记忆库内移动笔记。"""
        source = _safe_note_path(notes, note)
        destination = _safe_note_path(
            notes,
            f"{str(destination_folder).strip('/')}/{source.stem}",
            allow_missing=True,
        )
        if destination.exists():
            raise ValueError("目标位置已有同名笔记，拒绝覆盖。")
        destination.parent.mkdir(parents=True, exist_ok=True)
        source.replace(destination)
        return _note_payload(destination, notes, destination.read_text(encoding="utf-8"))

    return mcp


def _load_frame(path: Path, sheet_name: str = ""):
    import pandas as pd

    suffix = path.suffix.lower()
    if suffix == ".csv":
        return pd.read_csv(path)
    if suffix == ".tsv":
        return pd.read_csv(path, sep="\t")
    if suffix == ".json":
        return pd.read_json(path)
    return pd.read_excel(path, sheet_name=sheet_name or 0)


def _frame_records(frame, *, limit: int = 1_000) -> list[dict[str, Any]]:
    safe = frame.head(limit)
    return json.loads(safe.to_json(orient="records", force_ascii=False, date_format="iso"))


def _excel_sheet_names(path: Path) -> list[str]:
    if path.suffix.lower() in {".csv", ".tsv", ".json"}:
        return ["data"]
    import pandas as pd
    return list(pd.ExcelFile(path).sheet_names)


def build_excel(context: WorkspaceContext) -> FastMCP:
    mcp = FastMCP("ModelMirror Excel")

    @mcp.tool(annotations=READ_ONLY)
    def read_excel(file_id: FileId, sheet_name: str = "", max_rows: int = 200) -> dict[str, Any]:
        """读取受控表格文件，默认返回前 200 行，最多 1000 行。"""
        path = context.resolve_file(file_id)
        limit = max(1, min(int(max_rows), 1_000))
        frame = _load_frame(path, sheet_name)
        return {"filename": path.name, "rows": len(frame), "columns": list(map(str, frame.columns)), "data": _frame_records(frame, limit=limit), "truncated": len(frame) > limit}

    @mcp.tool(annotations=READ_ONLY)
    def get_excel_info(file_id: FileId) -> dict[str, Any]:
        """获取受控表格的大小、工作表和列信息。"""
        path = context.resolve_file(file_id)
        frame = _load_frame(path)
        return {"filename": path.name, "size_bytes": path.stat().st_size, "sheets": _excel_sheet_names(path), "rows": len(frame), "columns": list(map(str, frame.columns))}

    @mcp.tool(annotations=READ_ONLY)
    def get_sheet_names(file_id: FileId) -> dict[str, Any]:
        """列出受控表格中的工作表。"""
        path = context.resolve_file(file_id)
        return {"filename": path.name, "sheets": _excel_sheet_names(path)}

    @mcp.tool(annotations=READ_ONLY)
    def analyze_excel(file_id: FileId, sheet_name: str = "") -> dict[str, Any]:
        """对受控表格执行描述性统计。"""
        frame = _load_frame(context.resolve_file(file_id), sheet_name)
        summary = frame.describe(include="all").fillna("")
        return {"rows": len(frame), "columns": list(map(str, frame.columns)), "summary": json.loads(summary.to_json(force_ascii=False, date_format="iso"))}

    @mcp.tool(annotations=READ_ONLY)
    def filter_excel(file_id: FileId, column: str, operator: Literal["eq", "ne", "gt", "gte", "lt", "lte", "contains"], value: str, sheet_name: str = "", max_rows: int = 200) -> dict[str, Any]:
        """按单列条件筛选受控表格。"""
        frame = _load_frame(context.resolve_file(file_id), sheet_name)
        if column not in frame.columns:
            raise ValueError("筛选列不存在。")
        series = frame[column]
        if operator == "contains":
            mask = series.astype(str).str.contains(str(value), case=False, regex=False, na=False)
        elif operator in {"gt", "gte", "lt", "lte"}:
            numeric = series.astype(float)
            target = float(value)
            mask = {"gt": numeric > target, "gte": numeric >= target, "lt": numeric < target, "lte": numeric <= target}[operator]
        else:
            mask = series.astype(str) == str(value)
            if operator == "ne":
                mask = ~mask
        filtered = frame[mask]
        limit = max(1, min(int(max_rows), 1_000))
        return {"matched_rows": len(filtered), "data": _frame_records(filtered, limit=limit), "truncated": len(filtered) > limit}

    @mcp.tool(annotations=READ_ONLY)
    def pivot_table(file_id: FileId, index: str, values: str, aggfunc: Literal["sum", "mean", "count", "min", "max"] = "sum", columns: str = "", sheet_name: str = "") -> dict[str, Any]:
        """为受控表格生成受限透视摘要。"""
        frame = _load_frame(context.resolve_file(file_id), sheet_name)
        for name in (index, values):
            if name not in frame.columns:
                raise ValueError("透视表列不存在。")
        if columns and columns not in frame.columns:
            raise ValueError("透视表分组列不存在。")
        result = frame.pivot_table(index=index, values=values, columns=columns or None, aggfunc=aggfunc).reset_index()
        return {"data": _frame_records(result, limit=1_000), "rows": len(result)}

    @mcp.tool(annotations=READ_ONLY)
    def data_summary(file_id: FileId, sheet_name: str = "") -> dict[str, Any]:
        """返回受控表格的数据类型、缺失值和唯一值摘要。"""
        frame = _load_frame(context.resolve_file(file_id), sheet_name)
        return {"rows": len(frame), "columns": [{"name": str(name), "dtype": str(frame[name].dtype), "missing": int(frame[name].isna().sum()), "unique": int(frame[name].nunique(dropna=True))} for name in frame.columns]}

    @mcp.tool(annotations=ARTIFACT_CREATE)
    def export_chart(file_id: FileId, x_column: str, y_column: str, chart_type: Literal["line", "bar", "scatter", "histogram"] = "bar", artifact_name: ArtifactName = "chart.png", sheet_name: str = "") -> dict[str, Any]:
        """从受控表格生成 PNG 图表产物，不修改输入文件。"""
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        frame = _load_frame(context.resolve_file(file_id), sheet_name)
        if x_column not in frame.columns or y_column not in frame.columns:
            raise ValueError("图表列不存在。")
        target = context.artifact_path(artifact_name, ".png")
        figure, axis = plt.subplots(figsize=(8, 5))
        if chart_type == "line": axis.plot(frame[x_column], frame[y_column])
        elif chart_type == "scatter": axis.scatter(frame[x_column], frame[y_column])
        elif chart_type == "histogram": axis.hist(frame[y_column].dropna())
        else: axis.bar(frame[x_column], frame[y_column])
        axis.set_xlabel(x_column); axis.set_ylabel(y_column); figure.tight_layout()
        figure.savefig(target, dpi=120); plt.close(figure)
        return context.artifact_payload(target)

    @mcp.tool(annotations=STATE_WRITE)
    def write_excel(data: list[dict[str, Any]], artifact_name: ArtifactName = "output.xlsx", sheet_name: str = "Sheet1") -> dict[str, Any]:
        """经审批后把受限表格数据写为新的 XLSX 产物。"""
        import pandas as pd
        if len(data) > MAX_EXCEL_ROWS or any(len(row) > MAX_EXCEL_COLUMNS for row in data):
            raise ValueError("Excel 写入最多 10000 行、200 列。")
        target = context.artifact_path(artifact_name, ".xlsx")
        pd.DataFrame(data).to_excel(target, index=False, sheet_name=str(sheet_name)[:31] or "Sheet1")
        return context.artifact_payload(target)

    @mcp.tool(annotations=STATE_WRITE)
    def update_excel(file_id: FileId, updates: list[dict[str, Any]], artifact_name: ArtifactName = "updated.xlsx", sheet_name: str = "") -> dict[str, Any]:
        """经审批后在输入副本上更新单元格并生成新 XLSX，绝不覆盖源文件。"""
        import pandas as pd
        if len(updates) > 10_000:
            raise ValueError("单次最多更新 10000 个单元格。")
        frame = _load_frame(context.resolve_file(file_id), sheet_name)
        for update in updates:
            row = int(update.get("row", 0))
            column = str(update.get("column", ""))
            if row < 0 or row >= len(frame) or column not in frame.columns:
                raise ValueError("更新位置越过表格范围。")
            frame.at[row, column] = update.get("value")
        target = context.artifact_path(artifact_name, ".xlsx")
        pd.DataFrame(frame).to_excel(target, index=False, sheet_name=(sheet_name or "Sheet1")[:31])
        return context.artifact_payload(target)

    return mcp


def _git_repository(context: WorkspaceContext) -> Path:
    if (context.input_root / ".git").is_dir():
        return context.input_root
    candidates = [path.parent for path in context.input_root.glob("*/.git") if path.is_dir()]
    if len(candidates) != 1:
        raise ValueError("工作区必须包含且只能包含一个 Git 仓库。")
    return candidates[0]


def _git(context: WorkspaceContext, *args: str, max_chars: int = MAX_INLINE_CHARS) -> dict[str, Any]:
    repository = _git_repository(context)
    command = [
        "git", "--no-optional-locks",
        "-c", "core.hooksPath=/dev/null",
        "-c", "core.fsmonitor=false",
        "-c", "credential.helper=",
        "-c", "protocol.file.allow=never",
        "-c", f"safe.directory={repository}",
        *args,
    ]
    env = {
        "PATH": "/usr/bin:/bin",
        "HOME": os.getenv("TMPDIR", "/tmp"),
        "LANG": "C.UTF-8",
        "LC_ALL": "C.UTF-8",
        "GIT_CONFIG_NOSYSTEM": "1",
        "GIT_CONFIG_GLOBAL": "/dev/null",
        "GIT_OPTIONAL_LOCKS": "0",
        "GIT_EXTERNAL_DIFF": "",
        "GIT_TERMINAL_PROMPT": "0",
    }
    result = subprocess.run(command, cwd=repository, env=env, stdin=subprocess.DEVNULL, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=30, check=False)
    output = result.stdout.decode("utf-8", errors="replace")
    error = result.stderr.decode("utf-8", errors="replace")
    if result.returncode != 0:
        raise ValueError((error or "Git 只读操作失败。")[:1_000])
    return {"output": output[:max_chars], "truncated": len(output) > max_chars}


def build_git(context: WorkspaceContext) -> FastMCP:
    mcp = FastMCP("ModelMirror Git Read Only")

    @mcp.tool(annotations=READ_ONLY)
    def git_status() -> dict[str, Any]:
        """查看上传仓库状态；禁止锁、Hook 和任何写入。"""
        return _git(context, "status", "--short", "--branch", "--untracked-files=all")

    @mcp.tool(annotations=READ_ONLY)
    def git_diff_unstaged() -> dict[str, Any]:
        """查看未暂存差异，禁用 external diff 和 textconv。"""
        return _git(context, "diff", "--no-ext-diff", "--no-textconv")

    @mcp.tool(annotations=READ_ONLY)
    def git_diff_staged() -> dict[str, Any]:
        """查看已暂存差异，禁用 external diff 和 textconv。"""
        return _git(context, "diff", "--cached", "--no-ext-diff", "--no-textconv")

    @mcp.tool(annotations=READ_ONLY)
    def git_diff(target: str) -> dict[str, Any]:
        """查看相对指定修订的只读差异。"""
        if not re.fullmatch(r"[A-Za-z0-9_./^-]{1,200}", str(target or "")) or ".." in str(target):
            raise ValueError("Git 修订标识无效。")
        return _git(context, "diff", "--no-ext-diff", "--no-textconv", str(target))

    @mcp.tool(annotations=READ_ONLY)
    def git_log(max_count: int = 30) -> dict[str, Any]:
        """查看最多 100 条提交历史。"""
        return _git(context, "log", f"--max-count={max(1, min(int(max_count), 100))}", "--date=iso-strict", "--pretty=format:%H%x09%ad%x09%an%x09%s")

    @mcp.tool(annotations=READ_ONLY)
    def git_show(revision: str = "HEAD") -> dict[str, Any]:
        """查看指定修订，禁用 external diff 和 textconv。"""
        if not re.fullmatch(r"[A-Za-z0-9_./^-]{1,200}", str(revision or "")) or ".." in str(revision):
            raise ValueError("Git 修订标识无效。")
        return _git(context, "show", "--no-ext-diff", "--no-textconv", "--format=fuller", str(revision))

    @mcp.tool(annotations=READ_ONLY)
    def git_branch() -> dict[str, Any]:
        """列出本地分支，不创建或切换分支。"""
        return _git(context, "branch", "--list", "--format=%(refname:short)%09%(objectname:short)%09%(subject)")

    return mcp


class OfficeDocumentParseError(ValueError):
    """Stable error raised without leaking a workspace path or source name."""

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code
        self.message = message


def _validate_office_handoff(context: WorkspaceContext, source: Path) -> str:
    """Authenticate one immutable main-process handoff before deep parsing."""

    invalid = OfficeDocumentParseError(
        "office_handoff_invalid",
        "The staged Office handoff metadata is invalid.",
    )
    marker_path = context.input_root / OFFICE_HANDOFF_MARKER_NAME
    try:
        if marker_path.is_symlink() or not marker_path.is_file():
            raise invalid
        flags = os.O_RDONLY
        if hasattr(os, "O_NOFOLLOW"):
            flags |= os.O_NOFOLLOW
        descriptor = os.open(marker_path, flags)
        try:
            marker_stat = os.fstat(descriptor)
            if (
                not stat.S_ISREG(marker_stat.st_mode)
                or marker_stat.st_size <= 0
                or marker_stat.st_size > OFFICE_HANDOFF_MARKER_MAX_BYTES
            ):
                raise invalid
            raw_marker = os.read(descriptor, OFFICE_HANDOFF_MARKER_MAX_BYTES + 1)
        finally:
            os.close(descriptor)
        if len(raw_marker) != marker_stat.st_size:
            raise invalid
        marker = json.loads(raw_marker.decode("utf-8"))
    except OfficeDocumentParseError:
        raise
    except (OSError, UnicodeError, json.JSONDecodeError, TypeError, ValueError) as exc:
        raise invalid from exc

    if not isinstance(marker, dict):
        raise invalid
    format_id = str(marker.get("format_id") or "").strip().casefold()
    source_name = str(marker.get("source_name") or "")
    expected_name = _OFFICE_FIXED_SOURCE_NAMES.get(format_id)
    expected_sha256 = str(marker.get("source_sha256") or "")
    if (
        marker.get("owner") != OFFICE_HANDOFF_MARKER_OWNER
        or marker.get("workspace_id") != context.workspace_id
        or expected_name is None
        or source_name != expected_name
        or source.name != expected_name
        or source.parent.resolve() != context.input_root.resolve()
        or _OFFICE_SHA256_PATTERN.fullmatch(expected_sha256) is None
    ):
        raise invalid

    try:
        flags = os.O_RDONLY
        if hasattr(os, "O_NOFOLLOW"):
            flags |= os.O_NOFOLLOW
        descriptor = os.open(source, flags)
        try:
            source_stat = os.fstat(descriptor)
            if (
                not stat.S_ISREG(source_stat.st_mode)
                or source_stat.st_size <= 0
                or source_stat.st_size > MAX_OFFICE_INPUT_BYTES
            ):
                raise invalid
            digest = hashlib.sha256()
            received = 0
            while received <= MAX_OFFICE_INPUT_BYTES:
                chunk = os.read(
                    descriptor,
                    min(1024 * 1024, MAX_OFFICE_INPUT_BYTES + 1 - received),
                )
                if not chunk:
                    break
                digest.update(chunk)
                received += len(chunk)
        finally:
            os.close(descriptor)
    except OfficeDocumentParseError:
        raise
    except OSError as exc:
        raise invalid from exc
    if (
        received != source_stat.st_size
        or received > MAX_OFFICE_INPUT_BYTES
        or not hmac.compare_digest(digest.hexdigest(), expected_sha256)
    ):
        raise OfficeDocumentParseError(
            "office_handoff_integrity_failed",
            "The staged Office source failed its integrity check.",
        )
    return format_id


class _OfficeDocumentBuilder:
    def __init__(self) -> None:
        self.sections: list[dict[str, Any]] = []
        self.warnings: list[str] = []
        self.extracted_chars = 0
        self.retained_chars = 0
        self.truncated = False

    def warn(self, message: str) -> None:
        clean = _clean_office_text(message)[:500]
        if clean and clean not in self.warnings and len(self.warnings) < 20:
            self.warnings.append(clean)

    def add(self, text: str, **source: Any) -> None:
        clean = _clean_office_text(text)
        if not clean:
            return
        self.extracted_chars += len(clean)
        remaining = clean
        while remaining:
            capacity = MAX_OFFICE_EXTRACTED_CHARS - self.retained_chars
            if capacity <= 0 or len(self.sections) >= MAX_OFFICE_SECTIONS:
                self.truncated = True
                return
            take = min(len(remaining), capacity, MAX_OFFICE_SECTION_CHARS)
            piece = remaining[:take]
            remaining = remaining[take:]
            section: dict[str, Any] = {"text": piece}
            for key in ("slide", "heading_path"):
                value = source.get(key)
                if value is not None:
                    section[key] = value
            self.sections.append(section)
            self.retained_chars += len(piece)
            if remaining:
                self.truncated = True

    def finish(self, *, format_id: str, title: str | None) -> dict[str, Any]:
        if not self.sections:
            raise OfficeDocumentParseError(
                "office_has_no_readable_content",
                "The Office document has no readable text, table, note, or image placeholder.",
            )
        if self.truncated:
            self.warn(
                "Content exceeded the safe extraction limit and was truncated; source markers were retained."
            )
        result = {
            "format": format_id,
            "title": _clean_office_text(title or "")[:500] or None,
            "sections": self.sections,
            "warnings": self.warnings,
            "extracted_chars": self.extracted_chars,
            "truncated": self.truncated,
        }
        return _fit_office_wire_payload(result)


def _clean_office_text(value: Any) -> str:
    text = str(value or "").replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    return text.strip()


def _fit_office_wire_payload(payload: dict[str, Any]) -> dict[str, Any]:
    """Keep the compact UTF-8 JSON result below the sidecar wire ceiling."""

    def size(value: dict[str, Any]) -> int:
        return len(
            json.dumps(value, ensure_ascii=False, separators=(",", ":")).encode(
                "utf-8"
            )
        )

    if size(payload) <= MAX_OFFICE_WIRE_BYTES:
        return payload

    payload["truncated"] = True
    warning = "The structured Office result exceeded 2 MiB and was safely truncated."
    if warning not in payload["warnings"]:
        payload["warnings"].append(warning)
    sections = payload["sections"]
    low, high = 1, len(sections)
    while low < high:
        middle = (low + high + 1) // 2
        candidate = {**payload, "sections": sections[:middle]}
        if size(candidate) <= MAX_OFFICE_WIRE_BYTES:
            low = middle
        else:
            high = middle - 1
    payload["sections"] = sections[:low]
    if size(payload) <= MAX_OFFICE_WIRE_BYTES:
        return payload

    # A single section is bounded to 20k characters, but metadata can still be
    # adversarial. Remove optional heading metadata before truncating text.
    payload["sections"][0].pop("heading_path", None)
    text = payload["sections"][0]["text"]
    low, high = 1, len(text)
    while low < high:
        middle = (low + high + 1) // 2
        payload["sections"][0]["text"] = text[:middle]
        if size(payload) <= MAX_OFFICE_WIRE_BYTES:
            low = middle
        else:
            high = middle - 1
    payload["sections"][0]["text"] = text[:low]
    if size(payload) > MAX_OFFICE_WIRE_BYTES:
        raise OfficeDocumentParseError(
            "office_result_too_large",
            "The Office extraction result could not be represented safely.",
        )
    return payload


def _safe_office_member_name(value: str, *, format_id: str) -> str:
    clean = str(value or "")
    if (
        not clean
        or len(clean) > 512
        or "\\" in clean
        or "\x00" in clean
        or clean.startswith("/")
        or ":" in clean.split("/", 1)[0]
    ):
        raise OfficeDocumentParseError(
            f"unsafe_{format_id}_container",
            "The Office package contains an unsafe internal path.",
        )
    path = PurePosixPath(clean)
    if any(part in {"", ".", ".."} or len(part) > 160 for part in path.parts):
        raise OfficeDocumentParseError(
            f"unsafe_{format_id}_container",
            "The Office package contains an unsafe internal path.",
        )
    return path.as_posix()


def _unsupported_office_member(name: str) -> bool:
    clean = f"/{name.casefold()}"
    allowed_printer_settings = bool(
        re.fullmatch(r"/ppt/printersettings/printersettings[0-9]+\.bin", clean)
    )
    return (
        clean.endswith("vbaproject.bin")
        or clean.endswith("vbadata.xml")
        or "/activex/" in clean
        or "/embeddings/" in clean
        or "/oleobjects/" in clean
        or (clean.endswith(".bin") and not allowed_printer_settings)
    )


def _read_office_member(
    archive: zipfile.ZipFile,
    entry: zipfile.ZipInfo,
    *,
    format_id: str,
) -> bytes:
    if entry.file_size > MAX_OFFICE_XML_BYTES:
        raise OfficeDocumentParseError(
            f"{format_id}_complexity_limit_exceeded",
            "The Office package exceeds the safe XML complexity limit.",
        )
    with archive.open(entry) as stream:
        content = stream.read(MAX_OFFICE_XML_BYTES + 1)
    if len(content) > MAX_OFFICE_XML_BYTES:
        raise OfficeDocumentParseError(
            f"{format_id}_complexity_limit_exceeded",
            "The Office package exceeds the safe XML complexity limit.",
        )
    return content


def _parse_bounded_office_xml(
    content: bytes,
    *,
    format_id: str,
    counters: dict[str, int],
) -> ElementTree.Element:
    try:
        from defusedxml import ElementTree as SafeElementTree

        root = SafeElementTree.fromstring(content)
    except Exception as exc:
        raise OfficeDocumentParseError(
            f"invalid_{format_id}",
            "The Office package contains invalid or unsafe XML.",
        ) from exc

    stack: list[tuple[ElementTree.Element, int]] = [(root, 1)]
    while stack:
        element, depth = stack.pop()
        counters["nodes"] += 1
        counters["attributes"] += len(element.attrib)
        counters["text"] += len(element.text or "") + len(element.tail or "")
        if (
            depth > MAX_OFFICE_XML_DEPTH
            or counters["nodes"] > MAX_OFFICE_XML_NODES
            or counters["attributes"] > MAX_OFFICE_XML_ATTRIBUTES
            or counters["text"] > MAX_OFFICE_XML_TEXT_CHARS
            or len(element.attrib) > 1_000
        ):
            raise OfficeDocumentParseError(
                f"{format_id}_complexity_limit_exceeded",
                "The Office XML structure exceeds the safe complexity limit.",
            )
        children = list(element)
        stack.extend((child, depth + 1) for child in children)
    return root


def _relationship_source_part(name: str, *, format_id: str) -> str:
    path = PurePosixPath(name)
    if path.name == ".rels" and path.parent.as_posix() == "_rels":
        return ""
    if path.parent.name != "_rels" or not path.name.endswith(".rels"):
        raise OfficeDocumentParseError(
            f"invalid_{format_id}",
            "The Office package contains an invalid relationship path.",
        )
    return (path.parent.parent / path.name[: -len(".rels")]).as_posix()


def _safe_inert_hyperlink(relationship_type: str, target: str) -> bool:
    if not relationship_type.casefold().rstrip("/").endswith("/hyperlink"):
        return False
    if (
        len(target) > 2_048
        or "\\" in target
        or any(ord(character) < 0x20 for character in target)
    ):
        return False
    try:
        parsed = urlsplit(target)
    except ValueError:
        return False
    scheme = parsed.scheme.casefold()
    if scheme in {"http", "https"}:
        return (
            bool(parsed.hostname)
            and parsed.username is None
            and parsed.password is None
        )
    return scheme == "mailto" and bool(parsed.path) and not parsed.netloc


def _resolve_office_target(
    source_part: str,
    target: str,
    *,
    format_id: str,
) -> str:
    try:
        clean = unquote(str(target or "").strip(), errors="strict")
    except (UnicodeError, ValueError) as exc:
        raise OfficeDocumentParseError(
            f"unsafe_{format_id}_container",
            "The Office package contains an unsafe relationship target.",
        ) from exc
    if (
        not clean
        or "\\" in clean
        or "\x00" in clean
        or "?" in clean
        or "#" in clean
        or ":" in clean.split("/", 1)[0]
        or any(ord(character) < 0x20 for character in clean)
    ):
        raise OfficeDocumentParseError(
            f"unsafe_{format_id}_container",
            "The Office package contains an unsafe relationship target.",
        )
    parts = [] if clean.startswith("/") or not source_part else list(
        PurePosixPath(source_part).parent.parts
    )
    for part in PurePosixPath(clean.lstrip("/")).parts:
        if part in {"", "."}:
            continue
        if part == "..":
            if not parts:
                raise OfficeDocumentParseError(
                    f"unsafe_{format_id}_container",
                    "The Office relationship escapes the package root.",
                )
            parts.pop()
            continue
        if len(part) > 160:
            raise OfficeDocumentParseError(
                f"unsafe_{format_id}_container",
                "The Office relationship target is too long.",
            )
        parts.append(part)
    if not parts:
        raise OfficeDocumentParseError(
            f"invalid_{format_id}",
            "The Office relationship target is empty.",
        )
    return PurePosixPath(*parts).as_posix()


def _validate_office_relationships(
    relationship_roots: list[tuple[str, ElementTree.Element]],
    *,
    format_id: str,
    available_names: set[str],
) -> None:
    for relationship_name, root in relationship_roots:
        source_part = _relationship_source_part(
            relationship_name, format_id=format_id
        )
        if source_part and source_part.casefold() not in available_names:
            raise OfficeDocumentParseError(
                f"invalid_{format_id}",
                "The Office relationship source is missing.",
            )
        seen_ids: set[str] = set()
        for element in root.iter():
            if element.tag.rsplit("}", 1)[-1] != "Relationship":
                continue
            relationship_id = str(element.attrib.get("Id") or "").strip()
            relationship_type = str(element.attrib.get("Type") or "").strip()
            target = str(element.attrib.get("Target") or "").strip()
            target_mode = str(element.attrib.get("TargetMode") or "").strip().casefold()
            if (
                not relationship_id
                or not relationship_type
                or not target
                or relationship_id in seen_ids
                or len(relationship_type) > 2_048
            ):
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office relationship definition is incomplete or duplicated.",
                )
            seen_ids.add(relationship_id)
            if target_mode == "external":
                if not _safe_inert_hyperlink(relationship_type, target):
                    raise OfficeDocumentParseError(
                        f"unsupported_{format_id}_feature",
                        "External Office resources are not supported; inert HTTPS, HTTP, and mailto hyperlinks are the only exception.",
                    )
                continue
            if target_mode:
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office relationship uses an unknown target mode.",
                )
            normalized_type = relationship_type.casefold().rstrip("/")
            if any(
                normalized_type.endswith(suffix)
                for suffix in _OFFICE_ACTIVE_RELATIONSHIP_SUFFIXES
            ):
                raise OfficeDocumentParseError(
                    f"unsupported_{format_id}_feature",
                    "Macros, ActiveX, OLE, embedded packages, and external content are not supported.",
                )
            resolved = _resolve_office_target(
                source_part, target, format_id=format_id
            )
            if resolved.casefold() not in available_names:
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office package references missing internal content.",
                )


def _deep_validate_office_package(data: bytes, *, format_id: str) -> dict[str, Any]:
    if format_id not in _OFFICE_MAIN_PARTS:
        raise OfficeDocumentParseError(
            "office_format_not_supported", "Only DOCX and PPTX are supported."
        )
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = archive.infolist()
            if not entries:
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}", "The Office package is empty."
                )
            if len(entries) > MAX_OFFICE_ENTRIES:
                raise OfficeDocumentParseError(
                    f"{format_id}_complexity_limit_exceeded",
                    "The Office package has too many internal entries.",
                )
            normalized: dict[str, zipfile.ZipInfo] = {}
            total_uncompressed = 0
            total_compressed = 0
            for entry in entries:
                name = _safe_office_member_name(entry.filename, format_id=format_id)
                folded = name.casefold()
                if folded in normalized:
                    raise OfficeDocumentParseError(
                        f"unsafe_{format_id}_container",
                        "The Office package contains duplicate internal paths.",
                    )
                normalized[folded] = entry
                if entry.flag_bits & 0x1:
                    raise OfficeDocumentParseError(
                        f"encrypted_{format_id}",
                        "Encrypted Office documents are not supported.",
                    )
                mode = (entry.external_attr >> 16) & 0o170000
                if mode not in {0, 0o040000, 0o100000}:
                    raise OfficeDocumentParseError(
                        f"unsafe_{format_id}_container",
                        "The Office package contains a link or special file.",
                    )
                if entry.compress_type not in _OFFICE_ALLOWED_COMPRESSION:
                    raise OfficeDocumentParseError(
                        f"unsafe_{format_id}_container",
                        "The Office package uses an unsupported compression method.",
                    )
                if entry.file_size > MAX_OFFICE_MEMBER_BYTES:
                    raise OfficeDocumentParseError(
                        f"{format_id}_complexity_limit_exceeded",
                        "An Office package member exceeds the safe size limit.",
                    )
                if (
                    entry.file_size >= 1024 * 1024
                    and entry.file_size
                    > max(1, entry.compress_size) * MAX_OFFICE_COMPRESSION_RATIO
                ):
                    raise OfficeDocumentParseError(
                        f"{format_id}_complexity_limit_exceeded",
                        "An Office package member exceeds the safe expansion ratio.",
                    )
                total_uncompressed += max(0, entry.file_size)
                total_compressed += max(0, entry.compress_size)
                if _unsupported_office_member(name):
                    raise OfficeDocumentParseError(
                        f"unsupported_{format_id}_feature",
                        "Macros, ActiveX, OLE, and embedded packages are not supported.",
                    )
            if (
                total_uncompressed > MAX_OFFICE_UNCOMPRESSED_BYTES
                or total_uncompressed
                > max(1, total_compressed) * MAX_OFFICE_COMPRESSION_RATIO
            ):
                raise OfficeDocumentParseError(
                    f"{format_id}_complexity_limit_exceeded",
                    "The Office package exceeds the safe expansion limit.",
                )

            main_part, expected_root, expected_content_type = _OFFICE_MAIN_PARTS[
                format_id
            ]
            required = {"[content_types].xml", main_part}
            if not required.issubset(normalized):
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office package is missing its required document structure.",
                )
            corrupt = archive.testzip()
            if corrupt is not None:
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office package failed its internal CRC check.",
                )

            counters = {"nodes": 0, "attributes": 0, "text": 0}
            relationship_roots: list[tuple[str, ElementTree.Element]] = []
            content_types_root: ElementTree.Element | None = None
            main_root: ElementTree.Element | None = None
            has_revisions = False
            for folded_name, entry in normalized.items():
                if not (folded_name.endswith(".xml") or folded_name.endswith(".rels")):
                    continue
                content = _read_office_member(
                    archive, entry, format_id=format_id
                )
                root = _parse_bounded_office_xml(
                    content, format_id=format_id, counters=counters
                )
                if folded_name == "[content_types].xml":
                    content_types_root = root
                if folded_name == main_part:
                    main_root = root
                if folded_name.endswith(".rels"):
                    if root.tag.rsplit("}", 1)[-1] != "Relationships":
                        raise OfficeDocumentParseError(
                            f"invalid_{format_id}",
                            "The Office relationship XML has an invalid root.",
                        )
                    relationship_roots.append((folded_name, root))
                for element in root.iter():
                    local_name = element.tag.rsplit("}", 1)[-1].casefold()
                    if local_name in {
                        "oleobject",
                        "object",
                        "control",
                        "altchunk",
                    }:
                        raise OfficeDocumentParseError(
                            f"unsupported_{format_id}_feature",
                            "The Office document contains an active or embedded object.",
                        )
                    if format_id == "docx" and local_name in {
                        "ins",
                        "del",
                        "movefrom",
                        "moveto",
                    }:
                        has_revisions = True
                    if local_name in {"instrtext", "fldsimple"}:
                        field_text = " ".join(
                            [element.text or "", *element.attrib.values()]
                        ).casefold()
                        if any(
                            token in field_text
                            for token in (
                                "includetext",
                                "includepicture",
                                "ddeauto",
                                "oleobject",
                            )
                        ):
                            raise OfficeDocumentParseError(
                                f"unsupported_{format_id}_feature",
                                "The Office document contains an active external-content field.",
                            )

            if content_types_root is None or main_root is None:
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}", "The Office package XML is incomplete."
                )
            if content_types_root.tag.rsplit("}", 1)[-1] != "Types":
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office package content types are invalid.",
                )
            if main_root.tag.rsplit("}", 1)[-1] != expected_root:
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office package content does not match its format.",
                )
            content_types = {
                str(element.attrib.get("PartName") or "").lstrip("/").casefold():
                str(element.attrib.get("ContentType") or "").casefold()
                for element in content_types_root.iter()
                if element.tag.rsplit("}", 1)[-1] == "Override"
            }
            if content_types.get(main_part) != expected_content_type:
                raise OfficeDocumentParseError(
                    f"invalid_{format_id}",
                    "The Office main content type does not match its format.",
                )
            if any(
                any(
                    token in content_type
                    for token in ("macroenabled", "vba", "activex", "oleobject")
                )
                for content_type in content_types.values()
            ):
                raise OfficeDocumentParseError(
                    f"unsupported_{format_id}_feature",
                    "The Office package declares macro, ActiveX, or OLE content.",
                )
            _validate_office_relationships(
                relationship_roots,
                format_id=format_id,
                available_names=set(normalized),
            )
            return {"has_revisions": has_revisions}
    except OfficeDocumentParseError:
        raise
    except (zipfile.BadZipFile, zipfile.LargeZipFile, OSError, EOFError) as exc:
        raise OfficeDocumentParseError(
            f"invalid_{format_id}",
            "The Office package is damaged or structurally invalid.",
        ) from exc
    except Exception as exc:
        raise OfficeDocumentParseError(
            f"invalid_{format_id}",
            "The Office package could not be validated safely.",
        ) from exc


def _docx_heading_level(paragraph: Any) -> int | None:
    style = getattr(paragraph, "style", None)
    values = (
        str(getattr(style, "name", "") or ""),
        str(getattr(style, "style_id", "") or ""),
    )
    for value in values:
        match = re.search(r"heading\s*([1-9])", value, flags=re.IGNORECASE)
        if match:
            return int(match.group(1))
    return None


def _docx_list_prefix(paragraph: Any) -> str:
    style = getattr(getattr(paragraph, "style", None), "name", "") or ""
    lowered = str(style).casefold()
    if "list number" in lowered:
        return "1. "
    if "list bullet" in lowered:
        return "• "
    try:
        if paragraph._p.pPr is not None and paragraph._p.pPr.numPr is not None:
            return "• "
    except (AttributeError, TypeError):
        pass
    return ""


def _render_inert_link(label: str, target: str) -> str:
    clean_label = _clean_office_text(label) or "link"
    clean_target = _clean_office_text(target)[:2_048]
    return (
        f"{clean_label} [link target: {clean_target}]"
        if clean_target
        else clean_label
    )


def _docx_paragraph_text(paragraph: Any) -> tuple[str, int]:
    parts: list[str] = []
    try:
        items = paragraph.iter_inner_content()
    except AttributeError:
        items = getattr(paragraph, "runs", ())
    for item in items:
        text = str(getattr(item, "text", "") or "")
        if item.__class__.__name__ == "Hyperlink":
            target = str(getattr(item, "url", "") or "")
            parts.append(_render_inert_link(text, target))
        else:
            parts.append(text)
    image_count = sum(
        1
        for node in paragraph._element.iter()
        if node.tag.rsplit("}", 1)[-1] in {"drawing", "pict"}
    )
    parts.extend("[image]" for _ in range(image_count))
    rendered = "".join(parts).strip()
    prefix = _docx_list_prefix(paragraph)
    return (prefix + rendered if rendered else ""), image_count


def _docx_table_text(table: Any, *, depth: int = 0) -> tuple[str, int]:
    if depth > 8:
        raise OfficeDocumentParseError(
            "docx_complexity_limit_exceeded",
            "The DOCX table nesting exceeds the safe depth limit.",
        )
    rows: list[str] = []
    images = 0
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            blocks: list[str] = []
            try:
                contents = cell.iter_inner_content()
            except AttributeError:
                contents = cell.paragraphs
            for item in contents:
                if item.__class__.__name__ == "Paragraph":
                    text, count = _docx_paragraph_text(item)
                    images += count
                    if text:
                        blocks.append(text)
                elif item.__class__.__name__ == "Table":
                    nested, count = _docx_table_text(item, depth=depth + 1)
                    images += count
                    if nested:
                        blocks.append("[nested table]\n" + nested)
            cells.append(" / ".join(blocks).replace("\t", " "))
        rows.append("\t".join(cells))
    return "[table]\n" + "\n".join(rows), images


def _extract_docx(data: bytes, validation: dict[str, Any]) -> dict[str, Any]:
    try:
        from docx import Document

        document = Document(io.BytesIO(data))
        builder = _OfficeDocumentBuilder()
        headings: list[str] = []
        image_count = 0
        for item in document.iter_inner_content():
            if item.__class__.__name__ == "Paragraph":
                text, count = _docx_paragraph_text(item)
                image_count += count
                if not text:
                    continue
                level = _docx_heading_level(item)
                if level is not None:
                    headings = headings[: level - 1]
                    while len(headings) < level - 1:
                        headings.append("")
                    headings.append(text[:200])
                source = {
                    "heading_path": [value for value in headings if value]
                    or None
                }
                builder.add(text, **source)
            elif item.__class__.__name__ == "Table":
                text, count = _docx_table_text(item)
                image_count += count
                builder.add(
                    text,
                    heading_path=[value for value in headings if value] or None,
                )
        if validation.get("has_revisions"):
            builder.warn(
                "Tracked revisions were detected; inserted, deleted, or moved revision content may not be extracted completely."
            )
        if image_count:
            builder.warn(
                "Document images are represented by inert placeholders; no vision model was called."
            )
        title = str(getattr(document.core_properties, "title", "") or "")
        if not title and headings:
            title = next((value for value in headings if value), "")
        return builder.finish(format_id="docx", title=title)
    except OfficeDocumentParseError:
        raise
    except Exception as exc:
        raise OfficeDocumentParseError(
            "docx_parse_failed", "The DOCX document could not be parsed safely."
        ) from exc


def _pptx_text_frame_text(text_frame: Any) -> str:
    paragraphs: list[str] = []
    for paragraph in text_frame.paragraphs:
        runs: list[str] = []
        for run in paragraph.runs:
            text = str(run.text or "")
            try:
                target = str(run.hyperlink.address or "")
            except (AttributeError, KeyError, ValueError):
                target = ""
            runs.append(_render_inert_link(text, target) if target else text)
        rendered = "".join(runs).strip()
        if rendered:
            paragraphs.append(rendered)
    return "\n".join(paragraphs)


def _pptx_table_text(table: Any) -> str:
    rows: list[str] = []
    for row in table.rows:
        cells = [
            _pptx_text_frame_text(cell.text_frame).replace("\t", " ")
            for cell in row.cells
        ]
        rows.append("\t".join(cells))
    return "[table]\n" + "\n".join(rows)


def _pptx_shape_blocks(shape: Any, *, depth: int = 0) -> tuple[list[str], int]:
    if depth > 8:
        raise OfficeDocumentParseError(
            "pptx_complexity_limit_exceeded",
            "The PPTX group nesting exceeds the safe depth limit.",
        )
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        blocks: list[str] = []
        images = 0
        for child in shape.shapes:
            child_blocks, child_images = _pptx_shape_blocks(child, depth=depth + 1)
            blocks.extend(child_blocks)
            images += child_images
        return blocks, images
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return ["[image]"], 1
    if getattr(shape, "has_table", False):
        return [_pptx_table_text(shape.table)], 0
    if getattr(shape, "has_chart", False):
        return ["[chart]"], 0
    if getattr(shape, "has_text_frame", False):
        text = _pptx_text_frame_text(shape.text_frame)
        return ([text] if text else []), 0
    if shape.shape_type == MSO_SHAPE_TYPE.GRAPHIC_FRAME:
        return ["[graphic object]"], 0
    return [], 0


def _extract_pptx(data: bytes) -> dict[str, Any]:
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        builder = _OfficeDocumentBuilder()
        first_slide_title = ""
        image_count = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_title = ""
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape is not None and getattr(title_shape, "has_text_frame", False):
                slide_title = _pptx_text_frame_text(title_shape.text_frame)
            if not first_slide_title and slide_title:
                first_slide_title = slide_title
            blocks: list[str] = []
            for shape in slide.shapes:
                shape_blocks, shape_images = _pptx_shape_blocks(shape)
                blocks.extend(shape_blocks)
                image_count += shape_images
            if bool(getattr(slide, "has_notes_slide", False)):
                notes_frame = slide.notes_slide.notes_text_frame
                notes = _clean_office_text(
                    getattr(notes_frame, "text", "") if notes_frame is not None else ""
                )
                if notes:
                    blocks.append("[speaker notes]\n" + notes)
            text = "\n\n".join(block for block in blocks if _clean_office_text(block))
            builder.add(
                text,
                slide=slide_number,
                heading_path=[slide_title[:200]] if slide_title else None,
            )
        if image_count:
            builder.warn(
                "Slide images are represented by inert placeholders; no vision model was called."
            )
        title = str(getattr(presentation.core_properties, "title", "") or "")
        return builder.finish(
            format_id="pptx", title=title or first_slide_title
        )
    except OfficeDocumentParseError:
        raise
    except Exception as exc:
        raise OfficeDocumentParseError(
            "pptx_parse_failed", "The PPTX presentation could not be parsed safely."
        ) from exc


def extract_office_document_payload(path: Path) -> dict[str, Any]:
    """Validate deeply, then extract one DOCX/PPTX into ParsedDocument wire data."""

    suffix = path.suffix.casefold()
    format_id = {".docx": "docx", ".pptx": "pptx"}.get(suffix)
    if format_id is None:
        raise OfficeDocumentParseError(
            "office_format_not_supported", "Only DOCX and PPTX are supported."
        )
    try:
        if path.is_symlink() or not path.is_file():
            raise OSError("not a regular file")
        size = path.stat().st_size
        if size < 1:
            raise OfficeDocumentParseError(
                f"invalid_{format_id}", "The Office document is empty."
            )
        if size > MAX_OFFICE_INPUT_BYTES:
            raise OfficeDocumentParseError(
                "office_file_too_large",
                "The Office document exceeds the 10 MiB input limit.",
            )
        data = path.read_bytes()
    except OfficeDocumentParseError:
        raise
    except OSError as exc:
        raise OfficeDocumentParseError(
            "office_file_unavailable", "The Office document is unavailable."
        ) from exc
    validation = _deep_validate_office_package(data, format_id=format_id)
    return (
        _extract_docx(data, validation)
        if format_id == "docx"
        else _extract_pptx(data)
    )


class OutputRenderDocumentError(RuntimeError):
    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code
        self.message = message


def _validate_output_render_handoff(
    context: WorkspaceContext, path: Path
) -> dict[str, Any]:
    marker_path = context.input_root / OUTPUT_RENDER_MARKER_NAME
    try:
        if marker_path.is_symlink() or not marker_path.is_file():
            raise OSError("marker unavailable")
        raw_marker = marker_path.read_bytes()
        if len(raw_marker) > 4 * 1024:
            raise OSError("marker too large")
        marker = json.loads(raw_marker.decode("utf-8"))
        if (
            not isinstance(marker, dict)
            or marker.get("owner") != OUTPUT_RENDER_MARKER_OWNER
            or marker.get("workspace_id") != context.workspace_id
            or marker.get("source_name") != "spec.json"
            or path.name != "spec.json"
            or path.parent != context.input_root
            or marker.get("format_id") not in OUTPUT_RENDER_FORMATS
        ):
            raise OSError("marker mismatch")
        digest = str(marker.get("source_sha256") or "")
        if not re.fullmatch(r"[0-9a-f]{64}", digest):
            raise OSError("digest invalid")
        raw_spec = path.read_bytes()
        if not raw_spec or len(raw_spec) > MAX_OUTPUT_RENDER_SPEC_BYTES:
            raise OSError("spec size invalid")
        if not hmac.compare_digest(hashlib.sha256(raw_spec).hexdigest(), digest):
            raise OSError("spec digest mismatch")
        spec = json.loads(raw_spec.decode("utf-8"))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise OutputRenderDocumentError(
            "output_render_handoff_invalid",
            "The output render handoff failed its integrity check.",
        ) from exc
    if not isinstance(spec, dict) or spec.get("format_id") != marker.get("format_id"):
        raise OutputRenderDocumentError(
            "output_render_handoff_invalid",
            "The output render handoff failed its integrity check.",
        )
    return _validate_output_render_spec(spec)


def _validate_output_render_spec(spec: dict[str, Any]) -> dict[str, Any]:
    allowed_root = {
        "format_id", "filename", "title", "content", "rows", "blocks", "sheets", "slides"
    }
    if set(spec) - allowed_root:
        raise OutputRenderDocumentError(
            "output_render_spec_invalid", "The output render specification is invalid."
        )
    format_id = str(spec.get("format_id") or "")
    definition = OUTPUT_RENDER_FORMATS.get(format_id)
    filename = str(spec.get("filename") or "")
    if definition is None or Path(filename).name != filename or not filename.lower().endswith(definition[0]):
        raise OutputRenderDocumentError(
            "output_render_spec_invalid", "The output render specification is invalid."
        )
    blocks = spec.get("blocks") or []
    sheets = spec.get("sheets") or []
    slides = spec.get("slides") or []
    if not isinstance(blocks, list) or not isinstance(sheets, list) or not isinstance(slides, list):
        raise OutputRenderDocumentError(
            "output_render_spec_invalid", "The output render specification is invalid."
        )
    if format_id in {"pdf", "docx"} and (not blocks or len(blocks) > 10_000):
        raise OutputRenderDocumentError(
            "output_render_spec_invalid", "The document blocks are invalid."
        )
    if format_id == "xlsx" and (not sheets or len(sheets) > MAX_OUTPUT_RENDER_SHEETS):
        raise OutputRenderDocumentError(
            "output_render_spec_invalid", "The workbook sheet count is invalid."
        )
    if format_id == "pptx" and (not slides or len(slides) > MAX_OUTPUT_RENDER_SLIDES):
        raise OutputRenderDocumentError(
            "output_render_spec_invalid", "The presentation slide count is invalid."
        )
    cells = 0
    strings = 0
    for value in _walk_output_render_values(spec):
        if isinstance(value, str):
            strings += len(value)
        elif isinstance(value, float) and not math.isfinite(value):
            raise OutputRenderDocumentError(
                "output_render_spec_invalid", "The output contains a non-finite number."
            )
    for sheet in sheets:
        if not isinstance(sheet, dict) or set(sheet) - {"name", "rows"}:
            raise OutputRenderDocumentError(
                "output_render_spec_invalid", "The workbook sheet is invalid."
            )
        name = str(sheet.get("name") or "").strip()
        rows = sheet.get("rows") or []
        if not name or len(name) > 31 or any(character in name for character in "[]:*?/\\"):
            raise OutputRenderDocumentError(
                "output_render_spec_invalid", "The workbook sheet name is invalid."
            )
        if not isinstance(rows, list):
            raise OutputRenderDocumentError(
                "output_render_spec_invalid", "The workbook rows are invalid."
            )
        for row in rows:
            if not isinstance(row, list) or len(row) > MAX_OUTPUT_RENDER_COLUMNS:
                raise OutputRenderDocumentError(
                    "output_render_spec_invalid", "The workbook column limit was exceeded."
                )
            cells += len(row)
    if cells > MAX_OUTPUT_RENDER_CELLS or strings > MAX_OUTPUT_RENDER_CHARS:
        raise OutputRenderDocumentError(
            "output_render_spec_invalid", "The output render specification exceeds its safe limits."
        )
    return spec


def _walk_output_render_values(value: Any):
    if isinstance(value, dict):
        for item in value.values():
            yield from _walk_output_render_values(item)
    elif isinstance(value, list):
        for item in value:
            yield from _walk_output_render_values(item)
    else:
        yield value


def _safe_output_blocks(raw: Any) -> list[dict[str, Any]]:
    if not isinstance(raw, list):
        raise OutputRenderDocumentError("output_render_spec_invalid", "Document blocks are invalid.")
    result: list[dict[str, Any]] = []
    for block in raw:
        if not isinstance(block, dict) or set(block) - {"kind", "text", "level", "items", "rows"}:
            raise OutputRenderDocumentError("output_render_spec_invalid", "A document block is invalid.")
        kind = block.get("kind")
        if kind not in {"heading", "paragraph", "list", "table"}:
            raise OutputRenderDocumentError("output_render_spec_invalid", "A document block is invalid.")
        result.append(block)
    return result


def _neutralize_spreadsheet_text(value: Any) -> Any:
    if not isinstance(value, str):
        return value
    stripped = value.lstrip()
    if not stripped or stripped[0] not in "=+-@":
        return value
    if stripped[0] in "+-":
        try:
            numeric = float(stripped)
            if math.isfinite(numeric):
                return value
        except ValueError:
            pass
    return "'" + value


def _render_docx_output(spec: dict[str, Any], target: Path) -> list[str]:
    from docx import Document

    document = Document()
    title = str(spec.get("title") or "").strip()
    if title:
        document.core_properties.title = title
        document.add_heading(title, level=0)
    for block in _safe_output_blocks(spec.get("blocks")):
        kind = block["kind"]
        if kind == "heading":
            document.add_heading(str(block.get("text") or ""), level=int(block.get("level") or 1))
        elif kind == "paragraph":
            document.add_paragraph(str(block.get("text") or ""))
        elif kind == "list":
            for item in block.get("items") or []:
                document.add_paragraph(str(item), style="List Bullet")
        else:
            rows = block.get("rows") or []
            columns = max((len(row) for row in rows), default=0)
            if not columns or columns > MAX_OUTPUT_RENDER_COLUMNS:
                raise OutputRenderDocumentError("output_render_spec_invalid", "A document table is invalid.")
            table = document.add_table(rows=len(rows), cols=columns)
            for row_index, row in enumerate(rows):
                for column_index, value in enumerate(row):
                    table.cell(row_index, column_index).text = str(value)
    document.save(target)
    return []


def _render_xlsx_output(spec: dict[str, Any], target: Path) -> list[str]:
    from openpyxl import Workbook

    workbook = Workbook(write_only=True)
    neutralized = False
    for sheet_spec in spec.get("sheets") or []:
        sheet = workbook.create_sheet(str(sheet_spec["name"]))
        for row in sheet_spec.get("rows") or []:
            rendered = []
            for value in row:
                safe = _neutralize_spreadsheet_text(value)
                neutralized = neutralized or safe != value
                rendered.append(safe)
            sheet.append(rendered)
    workbook.save(target)
    return ["Spreadsheet-like formulas were neutralized as text."] if neutralized else []


def _render_pptx_output(spec: dict[str, Any], target: Path) -> list[str]:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slides._sldIdLst.clear()
    for slide_spec in spec.get("slides") or []:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(slide_spec.get("title") or "")
        top = Inches(1.2)
        for block in _safe_output_blocks(slide_spec.get("blocks") or []):
            if top >= Inches(6.8):
                break
            if block["kind"] == "table":
                rows = block.get("rows") or []
                columns = max((len(row) for row in rows), default=0)
                if not rows or not columns:
                    continue
                height = min(Inches(4.8), Inches(0.35 * len(rows) + 0.2))
                table = slide.shapes.add_table(len(rows), columns, Inches(0.7), top, Inches(8.6), height).table
                for row_index, row in enumerate(rows):
                    for column_index, value in enumerate(row):
                        table.cell(row_index, column_index).text = str(value)
                top += height + Inches(0.15)
                continue
            text = str(block.get("text") or "")
            if block["kind"] == "list":
                text = "\n".join(f"• {item}" for item in block.get("items") or [])
            box = slide.shapes.add_textbox(Inches(0.8), top, Inches(8.4), Inches(0.8))
            box.text_frame.text = text
            top += Inches(0.9)
    presentation.save(target)
    return []


def _render_pdf_output(spec: dict[str, Any], target: Path) -> list[str]:
    from html import escape

    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import (
        ListFlowable,
        ListItem,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    font_name = "STSong-Light"
    pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    styles = getSampleStyleSheet()
    body = ParagraphStyle("ModelMirrorBody", parent=styles["BodyText"], fontName=font_name, fontSize=10, leading=15)
    heading = ParagraphStyle("ModelMirrorHeading", parent=styles["Heading1"], fontName=font_name, fontSize=16, leading=22)
    title_style = ParagraphStyle("ModelMirrorTitle", parent=heading, alignment=TA_CENTER, fontSize=20)
    story: list[Any] = []
    title = str(spec.get("title") or "").strip()
    if title:
        story.extend([Paragraph(escape(title), title_style), Spacer(1, 12)])
    for block in _safe_output_blocks(spec.get("blocks")):
        kind = block["kind"]
        if kind in {"heading", "paragraph"}:
            story.append(Paragraph(escape(str(block.get("text") or "")).replace("\n", "<br/>"), heading if kind == "heading" else body))
            story.append(Spacer(1, 6))
        elif kind == "list":
            story.append(ListFlowable([ListItem(Paragraph(escape(str(item)), body)) for item in block.get("items") or []], bulletType="bullet"))
            story.append(Spacer(1, 6))
        else:
            rows = [[Paragraph(escape(str(cell)), body) for cell in row] for row in block.get("rows") or []]
            if rows:
                table = Table(rows, repeatRows=1)
                table.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.25, colors.grey), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
                story.extend([table, Spacer(1, 6)])
    SimpleDocTemplate(str(target), pagesize=A4, title=title).build(story)
    return []


def render_output_document_payload(
    context: WorkspaceContext, path: Path
) -> dict[str, Any]:
    spec = _validate_output_render_handoff(context, path)
    format_id = str(spec["format_id"])
    suffix, media_type = OUTPUT_RENDER_FORMATS[format_id]
    target = context.artifact_path("output" + suffix, suffix)
    try:
        if format_id == "pdf":
            warnings = _render_pdf_output(spec, target)
        elif format_id == "docx":
            warnings = _render_docx_output(spec, target)
        elif format_id == "xlsx":
            warnings = _render_xlsx_output(spec, target)
        else:
            warnings = _render_pptx_output(spec, target)
        if target.is_symlink() or not target.is_file():
            raise OSError("render artifact unavailable")
        size = target.stat().st_size
        if size <= 0 or size > MAX_OUTPUT_RENDERED_BYTES:
            raise OSError("render artifact size invalid")
        payload = context.artifact_payload(target)
        payload.update(
            {"format_id": format_id, "media_type": media_type, "warnings": warnings}
        )
        return payload
    except OutputRenderDocumentError:
        raise
    except Exception as exc:
        target.unlink(missing_ok=True)
        raise OutputRenderDocumentError(
            "output_render_failed", "The file could not be rendered safely."
        ) from exc


def build_output_renderer(context: WorkspaceContext) -> FastMCP:
    mcp = FastMCP("ModelMirror Output Renderer")

    @mcp.tool(annotations=ARTIFACT_CREATE)
    def render_output_document(file_id: FileId) -> dict[str, Any]:
        """Render one bounded, structured file specification without network access."""

        try:
            path = context.resolve_file(file_id)
            return render_output_document_payload(context, path)
        except OutputRenderDocumentError as exc:
            raise ValueError(f"{exc.code}: {exc.message}") from None
        except Exception:
            raise ValueError(
                "output_render_failed: The file could not be rendered safely."
            ) from None

    return mcp


def build_office_parser(context: WorkspaceContext) -> FastMCP:
    mcp = FastMCP("ModelMirror Office Parser")

    @mcp.tool(annotations=READ_ONLY)
    def extract_office_document(file_id: FileId) -> dict[str, Any]:
        """Extract a deeply validated DOCX or PPTX selected by opaque file ID."""

        try:
            path = context.resolve_file(file_id)
            _validate_office_handoff(context, path)
            return extract_office_document_payload(path)
        except OfficeDocumentParseError as exc:
            raise ValueError(f"{exc.code}: {exc.message}") from None
        except Exception:
            raise ValueError(
                "office_parse_failed: The Office document could not be parsed safely."
            ) from None

    return mcp


def build_markitdown(context: WorkspaceContext) -> FastMCP:
    mcp = FastMCP("ModelMirror MarkItDown")

    @mcp.tool(annotations=ARTIFACT_CREATE)
    def convert_to_markdown(file_id: FileId, artifact_name: ArtifactName = "converted.md") -> dict[str, Any]:
        """把受控工作区文件转换为 Markdown；不接受 URL、URI 或宿主路径。"""
        from markitdown import MarkItDown
        path = context.resolve_file(file_id)
        result = MarkItDown(enable_plugins=False).convert_local(path)
        text = str(result.markdown or "")
        target = context.artifact_path(artifact_name, ".md")
        target.write_text(text, encoding="utf-8")
        payload = context.artifact_payload(target)
        payload.update({"preview": text[:20_000], "preview_truncated": len(text) > 20_000})
        return payload

    return mcp


BUILDERS = {
    "basic-memory-mcp": build_basic_memory,
    "excel-mcp-server": build_excel,
    "git-mcp": build_git,
    "markitdown-mcp": build_markitdown,
    "office-parser-mcp": build_office_parser,
    "output-renderer-mcp": build_output_renderer,
    **WAVE18A_BUILDERS,
    **WAVE18B_BUILDERS,
    **WAVE20_BUILDERS,
    **WAVE26_BUILDERS,
}

ADAPTER_TOOL_NAMES = {
    "basic-memory-mcp": (
        "read_note", "read_content", "view_note", "search_notes", "search",
        "fetch", "recent_activity", "list_directory", "build_context",
        "basic_memory_diagnostics", "write_note", "edit_note", "move_note",
    ),
    "excel-mcp-server": (
        "read_excel", "get_excel_info", "get_sheet_names", "analyze_excel",
        "filter_excel", "pivot_table", "data_summary", "export_chart",
        "write_excel", "update_excel",
    ),
    "git-mcp": (
        "git_status", "git_diff_unstaged", "git_diff_staged", "git_diff",
        "git_log", "git_show", "git_branch",
    ),
    "markitdown-mcp": ("convert_to_markdown",),
    "office-parser-mcp": ("extract_office_document",),
    "output-renderer-mcp": ("render_output_document",),
    **WAVE18A_TOOL_NAMES,
    **WAVE18B_TOOL_NAMES,
    **WAVE20_TOOL_NAMES,
    **WAVE26_TOOL_NAMES,
}


def main() -> int:
    parser = argparse.ArgumentParser(add_help=False)
    parser.add_argument("adapter_id", choices=sorted(BUILDERS))
    args = parser.parse_args()
    context = WorkspaceContext()
    BUILDERS[args.adapter_id](context).run(transport="stdio")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
