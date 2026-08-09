from __future__ import annotations

import asyncio
import base64
import hashlib
import io
import json
import os
import struct
import subprocess
import sys
import tempfile
import time
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace


try:
    from server.mcp.file_proxy import ALLOWED_ADAPTERS
    from server.sandbox_sidecar.file_mcp import (
        ADAPTER_TOOL_NAMES,
        BUILDERS,
        MAX_OFFICE_WIRE_BYTES,
        OFFICE_HANDOFF_MARKER_NAME,
        OFFICE_HANDOFF_MARKER_OWNER,
        OfficeDocumentParseError,
        _validate_office_handoff,
        extract_office_document_payload,
        opaque_file_id,
    )
except ModuleNotFoundError:
    from mcp.file_proxy import ALLOWED_ADAPTERS
    from sandbox_sidecar.file_mcp import (
        ADAPTER_TOOL_NAMES,
        BUILDERS,
        MAX_OFFICE_WIRE_BYTES,
        OFFICE_HANDOFF_MARKER_NAME,
        OFFICE_HANDOFF_MARKER_OWNER,
        OfficeDocumentParseError,
        _validate_office_handoff,
        extract_office_document_payload,
        opaque_file_id,
    )


_ONE_PIXEL_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


def _require_office_libraries() -> None:
    try:
        import docx  # noqa: F401
        import pptx  # noqa: F401
    except ModuleNotFoundError as exc:
        raise unittest.SkipTest("Office parser dependencies live in mcp-files") from exc


def _add_docx_hyperlink(paragraph, label: str, target: str) -> None:
    from docx.opc.constants import RELATIONSHIP_TYPE
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    relationship_id = paragraph.part.relate_to(
        target, RELATIONSHIP_TYPE.HYPERLINK, is_external=True
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)
    run = OxmlElement("w:r")
    text = OxmlElement("w:t")
    text.text = label
    run.append(text)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _make_docx(path: Path, *, long_text: str = "") -> None:
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.shared import Inches

    document = Document()
    document.core_properties.title = "Quarterly brief"
    document.add_heading("Overview", level=1)
    paragraph = document.add_paragraph("Read ")
    _add_docx_hyperlink(paragraph, "the source", "https://example.com/report")
    document.add_paragraph("First item", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "42"
    document.add_picture(io.BytesIO(_ONE_PIXEL_PNG), width=Inches(0.1))
    if long_text:
        document.add_paragraph(long_text)
    # Empty revision markup is enough to prove that the parser emits its
    # fidelity warning without manufacturing deleted text.
    document.element.body.append(OxmlElement("w:ins"))
    document.save(path)


def _make_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.core_properties.title = "Launch review"
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Milestone one"
    box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Reference"
    run.hyperlink.address = "https://example.com/launch"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "State"
    table.cell(1, 0).text = "Team"
    table.cell(1, 1).text = "Ready"
    slide.shapes.add_picture(io.BytesIO(_ONE_PIXEL_PNG), Inches(1), Inches(4), Inches(0.2), Inches(0.2))
    slide.notes_slide.notes_text_frame.text = "Confirm launch checklist."
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Milestone two"
    second.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1)).text = "Ship"
    presentation.save(path)


def _write_office_handoff(
    workspace: Path,
    source: Path,
    *,
    workspace_id: str,
    format_id: str,
) -> None:
    (workspace / OFFICE_HANDOFF_MARKER_NAME).write_text(
        json.dumps(
            {
                "owner": OFFICE_HANDOFF_MARKER_OWNER,
                "workspace_id": workspace_id,
                "created_at": time.time(),
                "source_name": source.name,
                "source_sha256": hashlib.sha256(source.read_bytes()).hexdigest(),
                "format_id": format_id,
            },
            separators=(",", ":"),
        ),
        encoding="utf-8",
    )


def _rewrite_zip(path: Path, transform) -> None:
    with zipfile.ZipFile(path, "r") as source:
        members = [(entry.filename, source.read(entry)) for entry in source.infolist()]
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as target:
        transformed = transform(dict(members))
        for name, content in transformed.items():
            target.writestr(name, content)
    path.write_bytes(output.getvalue())


def _add_relationship(path: Path, relationship_xml: str) -> None:
    def transform(members: dict[str, bytes]) -> dict[str, bytes]:
        name = "word/_rels/document.xml.rels"
        content = members[name].decode("utf-8")
        marker = "</Relationships>"
        members[name] = content.replace(marker, relationship_xml + marker).encode("utf-8")
        return members

    _rewrite_zip(path, transform)


class OfficeParserRegistrationTests(unittest.TestCase):
    def test_adapter_is_private_but_available_to_the_fixed_proxy(self) -> None:
        self.assertIn("office-parser-mcp", BUILDERS)
        self.assertEqual(
            ADAPTER_TOOL_NAMES["office-parser-mcp"],
            ("extract_office_document",),
        )
        self.assertIn("office-parser-mcp", ALLOWED_ADAPTERS)


class OfficeParserHandoffTests(unittest.TestCase):
    def _workspace(self, root: Path) -> tuple[SimpleNamespace, Path]:
        workspace_id = "mcpws_" + "9" * 32
        workspace = root / workspace_id
        workspace.mkdir()
        context = SimpleNamespace(
            workspace_id=workspace_id,
            input_root=workspace,
        )
        return context, workspace

    def test_valid_marker_binds_fixed_name_format_workspace_and_sha(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            context, workspace = self._workspace(Path(directory))
            source = workspace / "source.docx"
            source.write_bytes(b"validated-office-source")
            _write_office_handoff(
                workspace,
                source,
                workspace_id=context.workspace_id,
                format_id="docx",
            )
            self.assertEqual(_validate_office_handoff(context, source), "docx")

    def test_tampered_marker_and_format_are_rejected(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            context, workspace = self._workspace(Path(directory))
            source = workspace / "source.docx"
            source.write_bytes(b"validated-office-source")
            _write_office_handoff(
                workspace,
                source,
                workspace_id=context.workspace_id,
                format_id="docx",
            )
            marker_path = workspace / OFFICE_HANDOFF_MARKER_NAME
            marker = json.loads(marker_path.read_text(encoding="utf-8"))
            marker["owner"] = "another-subsystem"
            marker_path.write_text(json.dumps(marker), encoding="utf-8")
            with self.assertRaises(OfficeDocumentParseError) as caught:
                _validate_office_handoff(context, source)
            self.assertEqual(caught.exception.code, "office_handoff_invalid")

            marker["owner"] = OFFICE_HANDOFF_MARKER_OWNER
            marker["format_id"] = "pptx"
            marker["source_name"] = "source.pptx"
            marker_path.write_text(json.dumps(marker), encoding="utf-8")
            with self.assertRaises(OfficeDocumentParseError) as caught:
                _validate_office_handoff(context, source)
            self.assertEqual(caught.exception.code, "office_handoff_invalid")

    def test_tampered_source_fails_independent_sha_check(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            context, workspace = self._workspace(Path(directory))
            source = workspace / "source.pptx"
            source.write_bytes(b"validated-office-source")
            _write_office_handoff(
                workspace,
                source,
                workspace_id=context.workspace_id,
                format_id="pptx",
            )
            source.write_bytes(b"tampered-after-staging")
            with self.assertRaises(OfficeDocumentParseError) as caught:
                _validate_office_handoff(context, source)
            self.assertEqual(
                caught.exception.code,
                "office_handoff_integrity_failed",
            )

    def test_oversized_marker_is_rejected_before_json_decode(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            context, workspace = self._workspace(Path(directory))
            source = workspace / "source.docx"
            source.write_bytes(b"validated-office-source")
            (workspace / OFFICE_HANDOFF_MARKER_NAME).write_bytes(b"{" + b"x" * 4096)
            with self.assertRaises(OfficeDocumentParseError) as caught:
                _validate_office_handoff(context, source)
            self.assertEqual(caught.exception.code, "office_handoff_invalid")


class OfficeParserGoldenTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        _require_office_libraries()

    def test_docx_preserves_order_headings_lists_links_tables_and_images(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "golden.docx"
            _make_docx(path)
            payload = extract_office_document_payload(path)

        self.assertEqual(payload["format"], "docx")
        self.assertEqual(payload["title"], "Quarterly brief")
        text = "\n".join(section["text"] for section in payload["sections"])
        self.assertLess(text.index("Overview"), text.index("Read"))
        self.assertLess(text.index("Read"), text.index("[table]"))
        self.assertIn("• First item", text)
        self.assertIn("the source [link target: https://example.com/report]", text)
        self.assertIn("Metric\tValue", text)
        self.assertIn("[image]", text)
        self.assertTrue(any(section.get("heading_path") == ["Overview"] for section in payload["sections"]))
        self.assertTrue(any("Tracked revisions" in warning for warning in payload["warnings"]))
        self.assertTrue(any("no vision model" in warning for warning in payload["warnings"]))

    def test_pptx_preserves_slides_tables_notes_links_and_images(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "golden.pptx"
            _make_pptx(path)
            payload = extract_office_document_payload(path)

        self.assertEqual(payload["format"], "pptx")
        self.assertEqual(payload["title"], "Launch review")
        self.assertEqual([section["slide"] for section in payload["sections"]], [1, 2])
        first = payload["sections"][0]["text"]
        self.assertIn("Milestone one", first)
        self.assertIn("Reference [link target: https://example.com/launch]", first)
        self.assertIn("Owner\tState", first)
        self.assertIn("[speaker notes]\nConfirm launch checklist.", first)
        self.assertIn("[image]", first)
        self.assertIn("Milestone two", payload["sections"][1]["text"])
        self.assertTrue(any("no vision model" in warning for warning in payload["warnings"]))

    def test_output_is_bounded_by_character_and_utf8_wire_limits(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "large.docx"
            long_text = "".join(
                hashlib.sha256(str(index).encode("ascii")).hexdigest()
                for index in range(8_000)
            )
            _make_docx(path, long_text=long_text)
            payload = extract_office_document_payload(path)
        wire = json.dumps(payload, ensure_ascii=False, separators=(",", ":")).encode("utf-8")
        self.assertLessEqual(len(wire), MAX_OFFICE_WIRE_BYTES)
        self.assertTrue(payload["truncated"])
        self.assertGreaterEqual(payload["extracted_chars"], 500_000)


class OfficeParserSecurityTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        _require_office_libraries()

    def _assert_rejected(self, path: Path, expected_code: str) -> None:
        with self.assertRaises(OfficeDocumentParseError) as caught:
            extract_office_document_payload(path)
        self.assertEqual(caught.exception.code, expected_code)
        self.assertNotIn(str(path), caught.exception.message)
        self.assertNotIn(path.name, caught.exception.message)

    def test_deep_validation_rejects_malformed_ancillary_xml(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "bad-xml.docx"
            _make_docx(path)

            def transform(members: dict[str, bytes]) -> dict[str, bytes]:
                members["word/extra.xml"] = b"<broken"
                return members

            _rewrite_zip(path, transform)
            self._assert_rejected(path, "invalid_docx")

    def test_deep_validation_rejects_external_image_relationship(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "external.docx"
            _make_docx(path)
            _add_relationship(
                path,
                '<Relationship Id="rExternal" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="https://example.com/pixel.png" TargetMode="External"/>',
            )
            self._assert_rejected(path, "unsupported_docx_feature")

    def test_deep_validation_rejects_unsafe_external_hyperlink(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "unsafe-link.docx"
            _make_docx(path)
            _add_relationship(
                path,
                '<Relationship Id="rUnsafe" '
                'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
                'relationships/hyperlink" Target="javascript:alert(1)" '
                'TargetMode="External"/>',
            )
            self._assert_rejected(path, "unsupported_docx_feature")

    def test_deep_validation_rejects_missing_internal_target(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "missing.docx"
            _make_docx(path)
            _add_relationship(
                path,
                '<Relationship Id="rMissing" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/missing.png"/>',
            )
            self._assert_rejected(path, "invalid_docx")

    def test_deep_validation_rejects_ole_and_expansion_bomb(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            ole_path = Path(directory) / "ole.docx"
            _make_docx(ole_path)

            def add_ole(members: dict[str, bytes]) -> dict[str, bytes]:
                members["word/embeddings/oleObject1.bin"] = b"unsafe"
                return members

            _rewrite_zip(ole_path, add_ole)
            self._assert_rejected(ole_path, "unsupported_docx_feature")

            bomb_path = Path(directory) / "bomb.docx"
            _make_docx(bomb_path)

            def add_bomb(members: dict[str, bytes]) -> dict[str, bytes]:
                members["word/media/bomb.dat"] = b"0" * (2 * 1024 * 1024)
                return members

            _rewrite_zip(bomb_path, add_bomb)
            self._assert_rejected(bomb_path, "docx_complexity_limit_exceeded")

    def test_deep_validation_rejects_crc_corruption(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "crc.docx"
            _make_docx(path)
            raw = bytearray(path.read_bytes())
            with zipfile.ZipFile(io.BytesIO(raw)) as archive:
                entry = archive.getinfo("word/document.xml")
            offset = entry.header_offset
            filename_length, extra_length = struct.unpack_from("<HH", raw, offset + 26)
            payload_offset = offset + 30 + filename_length + extra_length
            raw[payload_offset + max(1, entry.compress_size // 2)] ^= 0xFF
            path.write_bytes(raw)
            self._assert_rejected(path, "invalid_docx")


class OfficeParserStdioTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        _require_office_libraries()
        try:
            import mcp  # noqa: F401
        except ModuleNotFoundError as exc:
            raise unittest.SkipTest("MCP SDK lives in mcp-files") from exc

    def test_fastmcp_returns_parsed_document_as_structured_content(self) -> None:
        async def exercise() -> None:
            from mcp import ClientSession, StdioServerParameters
            from mcp.client.stdio import stdio_client

            with tempfile.TemporaryDirectory() as directory:
                root = Path(directory)
                workspace_id = "mcpws_" + "a" * 32
                input_root = root / "inputs"
                workspace = input_root / workspace_id
                workspace.mkdir(parents=True)
                path = workspace / "source.docx"
                _make_docx(path)
                _write_office_handoff(
                    workspace,
                    path,
                    workspace_id=workspace_id,
                    format_id="docx",
                )
                file_id = opaque_file_id(workspace_id, "source.docx")
                environment = dict(os.environ)
                environment.update(
                    {
                        "MCP_FILE_WORKSPACE_ID": workspace_id,
                        "MCP_FILE_INPUT_ROOT": str(input_root),
                        "MCP_FILE_OUTPUT_ROOT": str(root / "outputs"),
                        "MCP_FILE_MEMORY_ROOT": str(root / "memory"),
                    }
                )
                parameters = StdioServerParameters(
                    command=sys.executable,
                    args=[
                        "-m",
                        "server.sandbox_sidecar.file_mcp",
                        "office-parser-mcp",
                    ],
                    env=environment,
                )
                async with stdio_client(parameters) as (reader, writer):
                    async with ClientSession(reader, writer) as session:
                        await session.initialize()
                        result = await session.call_tool(
                            "extract_office_document", {"file_id": file_id}
                        )
                self.assertFalse(result.isError)
                structured = result.structuredContent
                self.assertIsInstance(structured, dict)
                self.assertEqual(structured["format"], "docx")
                self.assertIn("sections", structured)
                self.assertEqual(json.loads(result.content[0].text), structured)

        asyncio.run(exercise())


@unittest.skipUnless(
    os.getenv("MODELMIRROR_TEST_FILE_SIDECAR") == "1",
    "requires the isolated mcp-files container",
)
class OfficeParserSidecarChainTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        _require_office_libraries()

    def test_socket_proxy_launches_the_landlocked_office_worker(self) -> None:
        async def exercise() -> None:
            from mcp import ClientSession, StdioServerParameters
            from mcp.client.stdio import stdio_client

            workspace_id = "mcpws_" + "b" * 32
            input_root = Path("/inputs")
            workspace = input_root / workspace_id
            workspace.mkdir(parents=True, exist_ok=True)
            path = workspace / "source.docx"
            _make_docx(path)
            _write_office_handoff(
                workspace,
                path,
                workspace_id=workspace_id,
                format_id="docx",
            )
            file_id = opaque_file_id(workspace_id, "source.docx")
            socket_path = Path("/run/modelmirror-files-mcp/files-mcp.sock")
            server = subprocess.Popen(
                [sys.executable, "-m", "sandbox_sidecar.file_server"],
                env={
                    **os.environ,
                    "PYTHONPATH": "/opt/modelmirror",
                    "MCP_FILES_SOCKET_PATH": str(socket_path),
                    "MCP_FILE_INPUT_ROOT": str(input_root),
                    "MCP_FILE_OUTPUT_ROOT": "/outputs",
                    "MCP_FILE_MEMORY_ROOT": "/memory",
                },
                stdout=subprocess.DEVNULL,
                stderr=subprocess.PIPE,
                text=True,
            )
            try:
                deadline = time.monotonic() + 5
                while not socket_path.exists() and time.monotonic() < deadline:
                    if server.poll() is not None:
                        break
                    time.sleep(0.05)
                self.assertTrue(socket_path.exists(), "file sidecar socket did not start")
                parameters = StdioServerParameters(
                    command=sys.executable,
                    args=[
                        "/workspace/server/mcp/file_proxy.py",
                        "office-parser-mcp",
                    ],
                    env={
                        **os.environ,
                        "MCP_FILE_WORKSPACE_ID": workspace_id,
                        "MCP_FILES_SOCKET_PATH": str(socket_path),
                    },
                )
                async with stdio_client(parameters) as (reader, writer):
                    async with ClientSession(reader, writer) as session:
                        await session.initialize()
                        result = await asyncio.wait_for(
                            session.call_tool(
                                "extract_office_document", {"file_id": file_id}
                            ),
                            timeout=30,
                        )
                self.assertFalse(result.isError)
                self.assertEqual(result.structuredContent["format"], "docx")
            finally:
                server.terminate()
                try:
                    server.wait(timeout=3)
                except subprocess.TimeoutExpired:
                    server.kill()
                    server.wait(timeout=3)
                error = server.stderr.read() if server.stderr else ""
                if server.stderr:
                    server.stderr.close()
                if server.returncode not in {0, -15}:
                    self.fail(f"file sidecar exited unexpectedly: {error[:500]}")

        asyncio.run(exercise())


if __name__ == "__main__":
    unittest.main(verbosity=2)
